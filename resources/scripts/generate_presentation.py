#!/usr/bin/env python3
"""
Generate SBOM / Dependency-Track presentation (.pptx) and Confluence page.
Usage:
    pip install python-pptx pillow requests
    python3 generate_presentation.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from PIL import Image as _PILImage
import os

# ── Screenshot helper ──────────────────────────────────────────────────────
# generate_docs.sh downloads real DT screenshots to /screenshots before
# running this script. If the file is present we embed it; otherwise the
# slide falls back to the detailed simulated UI already in the deck.
_SCREENSHOT_DIR = os.environ.get('SCREENSHOT_DIR', '/screenshots')

def _screenshot(name: str):
    """Return absolute path to a downloaded DT screenshot, or None."""
    path = os.path.join(_SCREENSHOT_DIR, name)
    return path if os.path.isfile(path) else None

# ── Brand colours ──────────────────────────────────────────────────────────
NAVY      = RGBColor(0x1e, 0x3a, 0x8a)   # deep navy
DARK_NAVY = RGBColor(0x0c, 0x1a, 0x44)   # deeper navy (backgrounds/base layers)
TEAL      = RGBColor(0x06, 0x74, 0x7c)   # dependency-track teal
DARK_TEAL = RGBColor(0x04, 0x51, 0x5a)   # darker teal (decorative)
AMBER     = RGBColor(0xd9, 0x7a, 0x06)   # warning
RED       = RGBColor(0xb9, 0x1c, 0x1c)   # critical
GREEN     = RGBColor(0x15, 0x80, 0x3d)   # ok / low
WHITE     = RGBColor(0xff, 0xff, 0xff)
GREY      = RGBColor(0x64, 0x74, 0x8b)   # muted text
LGREY     = RGBColor(0xf0, 0xf4, 0xf8)   # slide background
MID_GREY  = RGBColor(0xde, 0xe3, 0xea)   # card borders / dividers
BLACK     = RGBColor(0x0f, 0x17, 0x2a)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ── helpers ────────────────────────────────────────────────────────────────

def add_rect(slide, left, top, width, height, fill, alpha=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    return shape


def add_text_box(slide, text, left, top, width, height,
                 font_size=18, bold=False, color=BLACK, align=PP_ALIGN.LEFT,
                 italic=False, wrap=True, line_spacing=None):
    txb = slide.shapes.add_textbox(left, top, width, height)
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing is not None:
        p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def add_bullet_box(slide, items, left, top, width, height,
                   font_size=16, color=BLACK, title=None, title_size=18,
                   title_color=NAVY, bullet="▸ ", bold_first=False,
                   space_before=3, line_spacing=1.2):
    """Bullet list with consistent paragraph spacing and line height."""
    txb = slide.shapes.add_textbox(left, top, width, height)
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    if title:
        p = tf.paragraphs[0]
        p.space_after = Pt(6)
        run = p.add_run()
        run.text = title
        run.font.size = Pt(title_size)
        run.font.bold = True
        run.font.color.rgb = title_color
        first = False
    for i, item in enumerate(items):
        if first and i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(space_before)
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = f"{bullet}{item}"
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        if bold_first and i == 0:
            run.font.bold = True
    return txb


def nav_bar(slide, title, color=NAVY):
    """Top nav bar with teal left and bottom accent strips (0.85″ total height)."""
    BAR_H = Inches(0.85)
    add_rect(slide, 0, 0, SLIDE_W, BAR_H, color)
    # Left teal accent strip
    add_rect(slide, 0, 0, Inches(0.055), BAR_H, TEAL)
    # Bottom teal accent strip (inside the bar, flush to base)
    add_rect(slide, 0, BAR_H - Inches(0.05), SLIDE_W, Inches(0.05), TEAL)
    add_text_box(slide, title,
                 Inches(0.44), Inches(0.1), Inches(12.4), Inches(0.68),
                 font_size=26, bold=True, color=WHITE)


def footer_bar(slide, text="SBOM & Vulnerability Management  ·  Dependency-Track"):
    add_rect(slide, 0, SLIDE_H - Inches(0.4), SLIDE_W, Inches(0.02), TEAL)   # thin separator
    add_rect(slide, 0, SLIDE_H - Inches(0.38), SLIDE_W, Inches(0.38), NAVY)
    add_text_box(slide, text,
                 Inches(0.35), SLIDE_H - Inches(0.36), Inches(12.5), Inches(0.34),
                 font_size=10, color=RGBColor(0xcb, 0xd5, 0xe1), align=PP_ALIGN.CENTER)


def chip(slide, text, left, top, width, height, bg, fg=WHITE, font_size=13):
    add_rect(slide, left, top, width, height, bg)
    add_text_box(slide, text, left, top, width, height,
                 font_size=font_size, bold=True, color=fg, align=PP_ALIGN.CENTER)


# ── Slide builders ─────────────────────────────────────────────────────────

def slide_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # ── Background: dark navy base + lighter navy panel ──────────────────
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, DARK_NAVY)          # very dark base
    add_rect(slide, 0, 0, SLIDE_W, Inches(5.08), NAVY)           # lighter navy top section

    # ── Decorative geometry — top-right corner accent ────────────────────
    add_rect(slide, Inches(9.8),  0, Inches(3.53), Inches(3.5), DARK_TEAL)   # teal block
    add_rect(slide, Inches(11.2), 0, Inches(2.13), Inches(1.8), TEAL)        # brighter teal corner

    # ── Teal divider bar between sections ────────────────────────────────
    add_rect(slide, 0, Inches(5.08), SLIDE_W, Inches(0.1), TEAL)

    # ── Main content ─────────────────────────────────────────────────────
    add_text_box(slide, "SBOM & Vulnerability Management",
                 Inches(0.75), Inches(1.15), Inches(9.5), Inches(1.2),
                 font_size=46, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    add_text_box(slide,
                 "Automating Software Composition Analysis with Dependency-Track",
                 Inches(0.75), Inches(2.52), Inches(9.5), Inches(0.72),
                 font_size=22, color=RGBColor(0xc4, 0xce, 0xe0), align=PP_ALIGN.LEFT)
    add_text_box(slide, "Know every component. Manage every vulnerability.",
                 Inches(0.75), Inches(3.38), Inches(9.5), Inches(0.55),
                 font_size=17, italic=True, color=TEAL, align=PP_ALIGN.LEFT)

    # ── Badge chips on dark bottom section ───────────────────────────────
    badge_data = [
        ("Automated SBOM Generation", TEAL),
        ("Continuous Monitoring",     RGBColor(0x1d, 0x4e, 0xd8)),
        ("Vulnerability Reporting",   RGBColor(0x7c, 0x3a, 0xed)),
        ("Maven Plugin Integration",  GREEN),
    ]
    for i, (label, col) in enumerate(badge_data):
        bx = Inches(0.55 + i * 3.1)
        by = Inches(5.4)
        bw = Inches(2.88)
        bh = Inches(0.64)
        add_rect(slide, bx, by, bw, bh, col)
        add_text_box(slide, label, bx, by, bw, bh,
                     font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    footer_bar(slide)


def slide_what_is_sbom(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, LGREY)
    nav_bar(slide, "What is an SBOM?")
    footer_bar(slide)

    # Left column — definition (white card with teal left accent)
    add_rect(slide, Inches(0.3), Inches(1.05), Inches(5.8), Inches(5.8), WHITE)
    add_rect(slide, Inches(0.3), Inches(1.05), Inches(0.055), Inches(5.8), TEAL)  # accent strip
    add_text_box(slide, "Software Bill of Materials",
                 Inches(0.48), Inches(1.18), Inches(5.5), Inches(0.5),
                 font_size=20, bold=True, color=NAVY)
    add_text_box(
        slide,
        "An SBOM is a formal, machine-readable inventory of every open-source "
        "and third-party component included in an application — along with their "
        "versions, licences, and supply-chain relationships.\n\n"
        "It is the foundation for understanding what your software is made of "
        "and identifying where security risks exist.",
        Inches(0.48), Inches(1.78), Inches(5.48), Inches(2.45),
        font_size=15, color=GREY, line_spacing=1.25
    )
    add_bullet_box(slide,
        ["Component name & version",
         "Package supplier / origin",
         "Dependency relationships",
         "Licence identifiers (SPDX)",
         "Known vulnerability references (CVEs)"],
        Inches(0.48), Inches(4.35), Inches(5.48), Inches(2.4),
        font_size=14, color=BLACK,
        title="An SBOM records:", title_size=15, title_color=TEAL
    )

    # Right column — CycloneDX callout (taller header strip)
    add_rect(slide, Inches(6.4), Inches(1.05), Inches(6.6), Inches(2.55), NAVY)
    add_text_box(slide, "CycloneDX Format",
                 Inches(6.55), Inches(1.18), Inches(6.3), Inches(0.48),
                 font_size=20, bold=True, color=WHITE)
    add_text_box(slide,
        "CycloneDX is an OWASP standard for SBOMs. It is natively supported "
        "by Dependency-Track, cdxgen, and the CycloneDX Maven / Gradle plugins.",
        Inches(6.55), Inches(1.75), Inches(6.3), Inches(1.6),
        font_size=14, color=RGBColor(0xcb, 0xd5, 0xe1), line_spacing=1.25
    )

    add_rect(slide, Inches(6.4), Inches(3.78), Inches(6.6), Inches(3.07), TEAL)
    add_text_box(slide, "Why do SBOMs matter?",
                 Inches(6.55), Inches(3.91), Inches(6.3), Inches(0.48),
                 font_size=20, bold=True, color=WHITE)
    add_bullet_box(slide,
        ["Visibility into ALL transitive dependencies",
         "Rapid impact assessment when a new CVE drops",
         "Licence compliance & open-source governance",
         "Supply-chain security & audit trail"],
        Inches(6.55), Inches(4.48), Inches(6.3), Inches(2.0),
        font_size=14, color=WHITE, bullet="✓  "
    )


def slide_how_sboms_generated(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, LGREY)
    nav_bar(slide, "How SBOMs Are Generated")
    footer_bar(slide)

    add_text_box(slide, "Two complementary approaches — both feed directly into Dependency-Track",
                 Inches(0.35), Inches(0.95), Inches(12.6), Inches(0.4),
                 font_size=15, italic=True, color=GREY)

    # Card 1 — cdxgen (any language) with left accent strip
    add_rect(slide, Inches(0.3), Inches(1.5), Inches(6.1), Inches(5.3), WHITE)
    add_rect(slide, Inches(0.3), Inches(1.5), Inches(6.1), Inches(0.55), NAVY)
    add_rect(slide, Inches(0.3), Inches(2.05), Inches(0.05), Inches(4.75), NAVY)  # left accent
    add_text_box(slide, "cdxgen  —  Any Language / Repository",
                 Inches(0.44), Inches(1.57), Inches(5.86), Inches(0.45),
                 font_size=16, bold=True, color=WHITE)
    add_bullet_box(slide,
        ["Polyglot support: Java, Node.js, Python, .NET, Go, Ruby …",
         "Scans the source tree and lock files automatically",
         "Runs as a Jenkins pipeline step via the cdxgenRepo() shared library call",
         "Produces a CycloneDX JSON SBOM in one command",
         "No changes needed to the target repository",
         "Supports pinned versions for reproducible builds"],
        Inches(0.48), Inches(2.2), Inches(5.76), Inches(4.3),
        font_size=14, color=BLACK
    )

    # Card 2 — Maven plugin with left accent strip
    add_rect(slide, Inches(6.9), Inches(1.5), Inches(6.1), Inches(5.3), WHITE)
    add_rect(slide, Inches(6.9), Inches(1.5), Inches(6.1), Inches(0.55), TEAL)
    add_rect(slide, Inches(6.9), Inches(2.05), Inches(0.05), Inches(4.75), TEAL)  # left accent
    add_text_box(slide, "CycloneDX Maven Plugin  —  In-House Java Code",
                 Inches(7.04), Inches(1.57), Inches(5.86), Inches(0.45),
                 font_size=16, bold=True, color=WHITE)
    add_bullet_box(slide,
        ["Added directly to pom.xml — no extra tooling required",
         "Runs during the standard Maven build lifecycle",
         "Generates a precise SBOM reflecting the exact resolved dependency tree",
         "Uploads the SBOM to Dependency-Track automatically at build time",
         "Keeps vulnerability data current with every CI build",
         "Developers get instant feedback on new vulnerabilities introduced"],
        Inches(7.08), Inches(2.2), Inches(5.76), Inches(4.3),
        font_size=14, color=BLACK
    )

    # Arrow connector hint
    add_text_box(slide, "Both routes → uploadSBOM() → Dependency-Track",
                 Inches(0.3), Inches(6.95), Inches(12.7), Inches(0.35),
                 font_size=13, italic=True, color=TEAL, align=PP_ALIGN.CENTER)


def slide_maven_plugin(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, LGREY)
    nav_bar(slide, "Maven Plugin — In-House Code Integration")
    footer_bar(slide)

    add_text_box(slide,
        "Add the CycloneDX Maven plugin to your pom.xml and your SBOM is "
        "generated and published to Dependency-Track on every build.",
        Inches(0.35), Inches(0.95), Inches(12.6), Inches(0.45),
        font_size=15, italic=True, color=GREY)

    # pom.xml snippet box
    add_rect(slide, Inches(0.3), Inches(1.5), Inches(6.2), Inches(5.3),
             RGBColor(0x1e, 0x29, 0x3b))
    code = (
        "<plugin>\n"
        "  <groupId>org.cyclonedx</groupId>\n"
        "  <artifactId>\n"
        "    cyclonedx-maven-plugin\n"
        "  </artifactId>\n"
        "  <version>2.9.1</version>\n"
        "  <executions>\n"
        "    <execution>\n"
        "      <phase>package</phase>\n"
        "      <goals>\n"
        "        <goal>makeAggregateBom</goal>\n"
        "      </goals>\n"
        "    </execution>\n"
        "  </executions>\n"
        "  <configuration>\n"
        "    <outputFormat>json</outputFormat>\n"
        "    <outputName>sbom</outputName>\n"
        "  </configuration>\n"
        "</plugin>"
    )
    add_text_box(slide, "pom.xml",
                 Inches(0.45), Inches(1.6), Inches(5.8), Inches(0.35),
                 font_size=12, bold=True, color=TEAL)
    add_text_box(slide, code,
                 Inches(0.4), Inches(2.0), Inches(5.95), Inches(4.7),
                 font_size=12, color=RGBColor(0xe2, 0xe8, 0xf0))

    # Right column — what this gives you
    add_rect(slide, Inches(6.8), Inches(1.5), Inches(6.2), Inches(2.3), NAVY)
    add_text_box(slide, "What this gives you",
                 Inches(6.95), Inches(1.6), Inches(5.9), Inches(0.4),
                 font_size=17, bold=True, color=WHITE)
    add_bullet_box(slide,
        ["Zero manual steps — SBOM published on every mvn package",
         "Vulnerability status visible in DT the moment the build lands",
         "Policy violations surface in the CI pipeline, not weeks later",
         "Full transitive dependency tree — nothing hidden"],
        Inches(6.95), Inches(2.1), Inches(5.9), Inches(1.6),
        font_size=13, color=WHITE, bullet="✓  "
    )

    add_rect(slide, Inches(6.8), Inches(4.0), Inches(6.2), Inches(2.8), TEAL)
    add_text_box(slide, "Pipeline flow",
                 Inches(6.95), Inches(4.1), Inches(5.9), Inches(0.4),
                 font_size=17, bold=True, color=WHITE)
    steps = [
        "1.  Developer pushes code",
        "2.  Jenkins triggers Maven build",
        "3.  cyclonedx-maven-plugin generates sbom.json",
        "4.  uploadSBOM() posts SBOM to Dependency-Track",
        "5.  DT enriches with NVD / OSV / GitHub Advisory data",
        "6.  Vulnerabilities visible immediately in DT dashboard",
    ]
    add_bullet_box(slide, steps,
                   Inches(6.95), Inches(4.6), Inches(5.9), Inches(2.1),
                   font_size=12, color=WHITE, bullet="")


def slide_detection_tools(prs):
    """SBOM generation tools + vulnerability intelligence sources reference slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, LGREY)
    nav_bar(slide, "Detection Tools & Vulnerability Intelligence Sources")
    footer_bar(slide)

    add_text_box(slide,
        "The pipeline combines best-in-class SBOM generators with Dependency-Track's "
        "multi-source vulnerability intelligence to achieve comprehensive coverage.",
        Inches(0.35), Inches(0.95), Inches(12.6), Inches(0.38),
        font_size=14, italic=True, color=GREY)

    # ── Left: SBOM Generation Tools ────────────────────────────────────────
    add_rect(slide, Inches(0.3), Inches(1.48), Inches(6.15), Inches(5.75), WHITE)
    add_rect(slide, Inches(0.3), Inches(1.48), Inches(6.15), Inches(0.52), NAVY)   # taller header
    add_text_box(slide, "SBOM Generation Tools",
                 Inches(0.44), Inches(1.55), Inches(5.87), Inches(0.4),
                 font_size=15, bold=True, color=WHITE)

    tools = [
        (NAVY,  "cdxgen",
         "Polyglot — Java, Node.js, Python, Go, .NET, Ruby & more.\n"
         "Runs as a Docker container; no host install required."),
        (TEAL,  "CycloneDX Maven Plugin",
         "Embedded in pom.xml; runs at mvn package phase.\n"
         "Generates a precise resolved-dependency-tree SBOM."),
        (RGBColor(0x7c, 0x3a, 0xed), "Syft  (Anchore)",
         "Scans Docker image layers to produce a container SBOM.\n"
         "Used via the syftImage() Jenkins shared-library step."),
        (AMBER, "Grype  (Anchore)",
         "Vulnerability scanner that consumes a Syft/cdxgen SBOM.\n"
         "Used via the grypeScan() step for pre-upload triage."),
        (GREEN, "CycloneDX Gradle / Node / .NET",
         "Ecosystem-specific plugins that mirror the Maven plugin's\n"
         "approach for Gradle, npm, and .NET SDK projects."),
    ]

    for i, (col, name, desc) in enumerate(tools):
        y = Inches(2.12 + i * 0.98)
        add_rect(slide, Inches(0.44), y, Inches(0.18), Inches(0.78), col)
        add_text_box(slide, name,
                     Inches(0.72), y + Inches(0.04), Inches(5.59), Inches(0.28),
                     font_size=12, bold=True, color=col)
        add_text_box(slide, desc,
                     Inches(0.72), y + Inches(0.33), Inches(5.59), Inches(0.47),
                     font_size=11, color=GREY, line_spacing=1.2)

    # ── Right: Vulnerability Intelligence Sources ──────────────────────────
    add_rect(slide, Inches(6.7), Inches(1.48), Inches(6.3), Inches(5.75), WHITE)
    add_rect(slide, Inches(6.7), Inches(1.48), Inches(6.3), Inches(0.52), TEAL)    # taller header
    add_text_box(slide, "Vulnerability Intelligence Sources",
                 Inches(6.84), Inches(1.55), Inches(6.02), Inches(0.4),
                 font_size=15, bold=True, color=WHITE)

    add_text_box(slide,
        "Dependency-Track queries these feeds automatically — no manual CVE searching:",
        Inches(6.84), Inches(2.12), Inches(6.02), Inches(0.32),
        font_size=12, color=BLACK)

    sources = [
        (RED,   "NVD — National Vulnerability Database",
         "NIST canonical CVE registry with CVSS v2/v3/v4 scores.\nThe primary authoritative source for all CVEs."),
        (AMBER, "OSV — Open Source Vulnerabilities",
         "Google's open database covering PyPI, npm, Maven, Go,\nRuby gems, Cargo, and more."),
        (NAVY,  "GitHub Advisory Database",
         "Curated advisories tied directly to GitHub-hosted packages.\nHighly accurate package-URL matching."),
        (GREEN, "Sonatype OSS Index",
         "Component intelligence focused on Maven Central and npm.\nIncludes licence and ecosystem data."),
        (RGBColor(0x7c, 0x3a, 0xed), "VulnDB  (optional commercial)",
         "Broader coverage than NVD with earlier disclosure.\nConfigurable as an additional feed in DT settings."),
    ]

    for i, (col, name, desc) in enumerate(sources):
        y = Inches(2.55 + i * 0.95)
        add_rect(slide, Inches(6.84), y, Inches(0.18), Inches(0.75), col)
        add_text_box(slide, name,
                     Inches(7.12), y + Inches(0.04), Inches(5.74), Inches(0.28),
                     font_size=12, bold=True, color=col)
        add_text_box(slide, desc,
                     Inches(7.12), y + Inches(0.33), Inches(5.74), Inches(0.44),
                     font_size=11, color=GREY, line_spacing=1.2)


def slide_dt_overview(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, LGREY)
    nav_bar(slide, "Dependency-Track — Overview")
    footer_bar(slide)

    # Intro
    add_text_box(slide,
        "Dependency-Track is an OWASP flagship platform that ingests SBOMs and "
        "continuously monitors every component for new vulnerabilities.",
        Inches(0.35), Inches(0.95), Inches(12.6), Inches(0.55),
        font_size=15, italic=True, color=GREY)

    # 4 feature cards
    cards = [
        (NAVY,  "Continuous Monitoring",
         "Checks every component against NVD, OSV, GitHub Advisory, and "
         "Sonatype OSS Index automatically — no manual CVE searches."),
        (TEAL,  "Policy Engine",
         "Define licence policies, severity thresholds, and component "
         "allow/block lists. Violations are surfaced instantly."),
        (RGBColor(0x7c, 0x3a, 0xed), "Portfolio View",
         "Group projects into collections. See aggregate risk across an "
         "entire product portfolio at a glance."),
        (GREEN, "REST API & Integrations",
         "Fully headless via REST API. Integrates with Jenkins, GitHub "
         "Actions, Slack notifications, and JIRA ticketing."),
    ]
    for i, (col, title, body) in enumerate(cards):
        x = Inches(0.3 + i * 3.26)
        add_rect(slide, x, Inches(1.7), Inches(3.1), Inches(1.6), col)
        add_text_box(slide, title,
                     x + Inches(0.12), Inches(1.8), Inches(2.86), Inches(0.45),
                     font_size=15, bold=True, color=WHITE)
        add_text_box(slide, body,
                     x + Inches(0.12), Inches(2.32), Inches(2.86), Inches(0.9),
                     font_size=12, color=RGBColor(0xe2, 0xe8, 0xf0))

    # Screenshot placeholder
    add_rect(slide, Inches(0.3), Inches(3.55), Inches(12.7), Inches(3.2),
             RGBColor(0xe2, 0xe8, 0xf0))
    add_rect(slide, Inches(0.3), Inches(3.55), Inches(12.7), Inches(0.38), NAVY)
    add_text_box(slide, "Dependency-Track — Project Portfolio Dashboard",
                 Inches(0.45), Inches(3.6), Inches(12.4), Inches(0.32),
                 font_size=13, bold=True, color=WHITE)

    # Simulated dashboard table headers
    headers = ["Project", "Version", "Critical", "High", "Medium", "Low", "Risk Score", "Last BOM"]
    col_w = [3.0, 1.4, 1.1, 1.1, 1.1, 1.1, 1.3, 1.6]
    x_pos = [0.3]
    for w in col_w[:-1]:
        x_pos.append(x_pos[-1] + w)

    for i, (h, w) in enumerate(zip(headers, col_w)):
        add_rect(slide, Inches(x_pos[i]), Inches(4.05),
                 Inches(w - 0.03), Inches(0.38), RGBColor(0x1e, 0x3a, 0x8a))
        add_text_box(slide, h,
                     Inches(x_pos[i] + 0.05), Inches(4.08),
                     Inches(w - 0.1), Inches(0.32),
                     font_size=11, bold=True, color=WHITE)

    rows = [
        ["my-billing-service",   "2.4.1", "0", "3", "7", "12", "4.2", "Today"],
        ["customer-api",         "1.8.0", "1", "5", "9", "4",  "7.8", "Today"],
        ["reporting-engine",     "3.1.2", "0", "1", "3", "8",  "2.1", "Yesterday"],
        ["auth-service",         "1.2.0", "2", "8", "11","6",  "9.1", "Today"],
    ]
    row_colors = [
        WHITE,
        RGBColor(0xff, 0xf1, 0xf2),
        WHITE,
        RGBColor(0xff, 0xf1, 0xf2),
    ]
    sev_colors = {
        "0": GREEN, "1": AMBER, "2": RED, "3": RED,
        "4": AMBER, "5": RED, "7": RED, "8": RED, "9": AMBER,
    }
    for r, (row, rc) in enumerate(zip(rows, row_colors)):
        for i, (val, w) in enumerate(zip(row, col_w)):
            add_rect(slide, Inches(x_pos[i]), Inches(4.5 + r * 0.42),
                     Inches(w - 0.03), Inches(0.4), rc)
            fc = sev_colors.get(val, BLACK) if i in (2, 3, 4, 5) else BLACK
            add_text_box(slide, val,
                         Inches(x_pos[i] + 0.05), Inches(4.53 + r * 0.42),
                         Inches(w - 0.1), Inches(0.35),
                         font_size=11, color=fc,
                         bold=(i in (2, 3) and val not in ("0",)))


def slide_vuln_workflow(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, LGREY)
    nav_bar(slide, "Managing Vulnerabilities — The Workflow")
    footer_bar(slide)

    add_text_box(slide,
        "Dependency-Track surfaces every finding. Your team then triages, "
        "analyses, and resolves each one — all recorded within the platform.",
        Inches(0.35), Inches(0.95), Inches(12.6), Inches(0.45),
        font_size=15, italic=True, color=GREY)

    # Flow steps
    steps = [
        (NAVY,  "1. Ingest SBOM",
         "SBOM uploaded via Jenkins pipeline. DT resolves all components "
         "and enriches with CVE data from multiple advisories databases."),
        (TEAL,  "2. Review Findings",
         "Open the project in DT. Findings are listed by severity — "
         "Critical, High, Medium, Low — with CVE details and CVSS scores."),
        (AMBER, "3. Triage",
         "Decide the appropriate action for each finding: Escalate to "
         "the engineering team, suppress, or mark as not affected."),
        (RGBColor(0x7c, 0x3a, 0xed), "4. Record Analysis",
         "Document your analysis inside DT: add comments, set the "
         "analysis state, and optionally suppress the finding."),
        (GREEN, "5. Remediate & Re-scan",
         "Engineering upgrades the dependency. The next pipeline run "
         "re-uploads the SBOM and the finding disappears automatically."),
    ]
    for i, (col, title, body) in enumerate(steps):
        x = Inches(0.25 + i * 2.57)
        add_rect(slide, x, Inches(1.6), Inches(2.45), Inches(0.52), col)
        add_text_box(slide, title,
                     x + Inches(0.08), Inches(1.68), Inches(2.3), Inches(0.38),
                     font_size=13, bold=True, color=WHITE)
        add_rect(slide, x, Inches(2.15), Inches(2.45), Inches(2.35), WHITE)
        add_text_box(slide, body,
                     x + Inches(0.1), Inches(2.25), Inches(2.28), Inches(2.2),
                     font_size=12, color=BLACK)

    # Findings screenshot placeholder
    add_rect(slide, Inches(0.25), Inches(4.65), Inches(12.75), Inches(2.08),
             RGBColor(0xe2, 0xe8, 0xf0))
    add_rect(slide, Inches(0.25), Inches(4.65), Inches(12.75), Inches(0.35), NAVY)
    add_text_box(slide,
                 "Dependency-Track — Findings tab for a project (screenshot)",
                 Inches(0.4), Inches(4.68), Inches(12.5), Inches(0.3),
                 font_size=12, bold=True, color=WHITE)

    # Simulated findings rows
    findings = [
        ("CVE-2021-44228", "log4j-core",         "2.14.1", "CRITICAL", "10.0", "Remote Code Execution via JNDI lookup"),
        ("CVE-2022-42003", "jackson-databind",    "2.13.0", "HIGH",     "7.5",  "Uncontrolled resource consumption"),
        ("CVE-2023-20863", "spring-expression",   "5.3.25", "MEDIUM",   "5.9",  "Regular expression denial of service"),
    ]
    f_colors = {"CRITICAL": RED, "HIGH": AMBER, "MEDIUM": RGBColor(0xca, 0x8a, 0x04), "LOW": GREEN}
    row_bg   = [WHITE, RGBColor(0xf8, 0xfa, 0xfc), WHITE]
    for r, (cve, comp, ver, sev, score, desc) in enumerate(findings):
        y = Inches(5.08 + r * 0.47)
        add_rect(slide, Inches(0.25), y, Inches(12.75), Inches(0.45), row_bg[r])
        vals = [cve, comp, ver, sev, score, desc]
        widths = [2.0, 2.2, 1.2, 1.2, 0.9, 5.15]
        x_off = [0.25]
        for w in widths[:-1]:
            x_off.append(x_off[-1] + w)
        for ci, (v, w) in enumerate(zip(vals, widths)):
            fc = f_colors.get(v, BLACK)
            add_text_box(slide, v, Inches(x_off[ci] + 0.05), y + Inches(0.04),
                         Inches(w - 0.1), Inches(0.38),
                         font_size=11, color=fc, bold=(ci == 3))


def slide_analysis_actions(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, LGREY)
    nav_bar(slide, "Vulnerability Analysis — Comments, Suppression & State")
    footer_bar(slide)

    add_text_box(slide,
        "Every finding in Dependency-Track can be analysed, commented on, "
        "and formally recorded — creating an audit trail of your security decisions.",
        Inches(0.35), Inches(0.95), Inches(12.6), Inches(0.45),
        font_size=15, italic=True, color=GREY)

    # Left — analysis panel mock
    add_rect(slide, Inches(0.3), Inches(1.55), Inches(7.1), Inches(5.65), WHITE)
    add_rect(slide, Inches(0.3), Inches(1.55), Inches(7.1), Inches(0.45), NAVY)
    add_text_box(slide, "Analysis Panel — CVE-2021-44228 / log4j-core 2.14.1",
                 Inches(0.42), Inches(1.6), Inches(6.9), Inches(0.38),
                 font_size=12, bold=True, color=WHITE)

    # State selector mock
    add_text_box(slide, "Analysis State",
                 Inches(0.45), Inches(2.12), Inches(2.5), Inches(0.32),
                 font_size=12, bold=True, color=NAVY)
    add_rect(slide, Inches(0.45), Inches(2.48), Inches(3.0), Inches(0.38),
             RGBColor(0xf1, 0xf5, 0xf9))
    add_text_box(slide, "▼  NOT_AFFECTED",
                 Inches(0.52), Inches(2.52), Inches(2.9), Inches(0.3),
                 font_size=12, color=GREEN, bold=True)

    # Justification
    add_text_box(slide, "Justification",
                 Inches(0.45), Inches(3.02), Inches(2.5), Inches(0.32),
                 font_size=12, bold=True, color=NAVY)
    add_rect(slide, Inches(0.45), Inches(3.38), Inches(3.0), Inches(0.38),
             RGBColor(0xf1, 0xf5, 0xf9))
    add_text_box(slide, "▼  PROTECTED_BY_MITIGATING_CONTROL",
                 Inches(0.52), Inches(3.42), Inches(2.9), Inches(0.3),
                 font_size=10, color=GREY)

    # Suppressed toggle
    add_text_box(slide, "Suppressed",
                 Inches(0.45), Inches(3.92), Inches(1.5), Inches(0.32),
                 font_size=12, bold=True, color=NAVY)
    add_rect(slide, Inches(1.85), Inches(3.94), Inches(0.65), Inches(0.28),
             GREEN)
    add_text_box(slide, "ON", Inches(1.85), Inches(3.94), Inches(0.65), Inches(0.28),
                 font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Comment box
    add_text_box(slide, "Analyst Comment",
                 Inches(0.45), Inches(4.37), Inches(2.5), Inches(0.32),
                 font_size=12, bold=True, color=NAVY)
    add_rect(slide, Inches(0.45), Inches(4.73), Inches(6.75), Inches(1.75),
             RGBColor(0xf8, 0xfa, 0xfc))
    add_text_box(slide,
        "Reviewed 2025-01-15 by J. Smith.\n"
        "Our deployment uses the JNDI lookup feature is disabled via "
        "log4j2.formatMsgNoLookups=true JVM flag set at container startup. "
        "Confirmed with DevOps — the attack vector is mitigated. "
        "Tracking upgrade to log4j 2.17.1 in Q1 sprint board (JIRA-4821).",
        Inches(0.55), Inches(4.83), Inches(6.6), Inches(1.55),
        font_size=11, color=BLACK)

    # Right — state legend with taller header
    add_rect(slide, Inches(7.65), Inches(1.55), Inches(5.35), Inches(5.65), WHITE)
    add_rect(slide, Inches(7.65), Inches(1.55), Inches(5.35), Inches(0.55), NAVY)   # header strip
    add_rect(slide, Inches(7.65), Inches(2.1), Inches(0.05), Inches(5.1), NAVY)      # left accent
    add_text_box(slide, "Available Analysis States",
                 Inches(7.8), Inches(1.63), Inches(5.1), Inches(0.42),
                 font_size=17, bold=True, color=WHITE)
    states = [
        (GREEN,  "NOT_AFFECTED",
         "The component is present but the vulnerability does not apply "
         "to this deployment context."),
        (AMBER,  "IN_TRIAGE",
         "The finding is being actively investigated. Interim state "
         "until a final decision is made."),
        (RED,    "EXPLOITABLE",
         "Confirmed — the vulnerability is present and exploitable. "
         "Immediate remediation required."),
        (GREY,   "RESOLVED",
         "The dependency has been upgraded and the vulnerability is no "
         "longer present in the component inventory."),
        (GREY,   "FALSE_POSITIVE",
         "The vulnerability does not apply — typically because the "
         "affected code path is never invoked."),
    ]
    for i, (col, name, desc) in enumerate(states):
        y = Inches(2.22 + i * 0.97)
        add_rect(slide, Inches(7.75), y, Inches(2.1), Inches(0.38), col)
        add_text_box(slide, name,
                     Inches(7.8), y + Inches(0.04), Inches(2.0), Inches(0.3),
                     font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text_box(slide, desc,
                     Inches(10.0), y, Inches(2.85), Inches(0.7),
                     font_size=11, color=BLACK, line_spacing=1.2)


def slide_suppression(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, LGREY)
    nav_bar(slide, "Suppression — Hiding Accepted Risk")
    footer_bar(slide)

    add_text_box(slide,
        "When a finding is accepted or not applicable, it can be suppressed "
        "so it no longer inflates your vulnerability counts — but it is never deleted.",
        Inches(0.35), Inches(0.95), Inches(12.6), Inches(0.45),
        font_size=15, italic=True, color=GREY)

    bullets_l = [
        "Suppressed findings are hidden from the default dashboard view",
        "They remain fully searchable and visible with the 'Show suppressed' filter",
        "The original CVE details, CVSS score, and analyst notes are preserved",
        "Suppression can be revoked at any time — full reversibility",
        "Suppressed findings are still exported in the vulnerability PDF report",
        "Audit trail records who suppressed a finding and when",
    ]
    add_rect(slide, Inches(0.3), Inches(1.55), Inches(6.2), Inches(4.5), WHITE)
    add_rect(slide, Inches(0.3), Inches(1.55), Inches(0.05), Inches(4.5), TEAL)   # left accent
    add_bullet_box(slide, bullets_l,
                   Inches(0.48), Inches(1.7), Inches(5.9), Inches(4.2),
                   font_size=14, color=BLACK,
                   title="How suppression works", title_size=16, title_color=NAVY)

    # Right — when to suppress
    add_rect(slide, Inches(6.8), Inches(1.55), Inches(6.2), Inches(4.5), NAVY)
    add_text_box(slide, "When should you suppress?",
                 Inches(6.95), Inches(1.68), Inches(5.9), Inches(0.42),
                 font_size=16, bold=True, color=WHITE)
    when = [
        "The vulnerability is in a library path that is never executed",
        "A mitigating control is already in place (WAF rule, JVM flag, etc.)",
        "You have accepted the risk and documented why",
        "The finding is a confirmed false positive",
        "A fix is not yet available and the risk is formally accepted",
    ]
    add_bullet_box(slide, when,
                   Inches(6.95), Inches(2.2), Inches(5.9), Inches(3.7),
                   font_size=14, color=WHITE, bullet="✓  ")

    # Warning strip at bottom
    add_rect(slide, Inches(0.3), Inches(6.18), Inches(12.7), Inches(0.52),
             RGBColor(0xff, 0xf7, 0xed))
    add_text_box(slide,
        "⚠  Suppression is not remediation — always pair suppression with a "
        "documented justification and a remediation timeline where applicable.",
        Inches(0.45), Inches(6.23), Inches(12.4), Inches(0.42),
        font_size=13, color=AMBER)


def slide_reporting(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, LGREY)
    nav_bar(slide, "Automated Reporting")
    footer_bar(slide)

    add_text_box(slide,
        "Jenkins pipelines generate shareable PDF reports from live "
        "Dependency-Track data — no manual export steps required.",
        Inches(0.35), Inches(0.95), Inches(12.6), Inches(0.45),
        font_size=15, italic=True, color=GREY)

    # Card 1 — Vulnerability report
    add_rect(slide, Inches(0.3), Inches(1.55), Inches(6.1), Inches(5.15), WHITE)
    add_rect(slide, Inches(0.3), Inches(1.55), Inches(6.1), Inches(0.55), RED)    # taller header
    add_rect(slide, Inches(0.3), Inches(2.1), Inches(0.05), Inches(4.6), RED)     # left accent
    add_text_box(slide, "Security Analysis (Vulnerability) Report",
                 Inches(0.44), Inches(1.62), Inches(5.82), Inches(0.42),
                 font_size=14, bold=True, color=WHITE)
    add_bullet_box(slide,
        ["One report per project or across all projects",
         "Executive summary: Critical / High / Medium / Low counts",
         "Overall risk level: CRITICAL / HIGH / MODERATE / LOW",
         "Per-vulnerability detail: CVE ID, component, version, CVSS score",
         "Analyst comments and suppression status included",
         "Landscape A4 PDF — ready to share with management or auditors",
         "Auto-archived as a Jenkins build artefact"],
        Inches(0.48), Inches(2.18), Inches(5.76), Inches(4.4),
        font_size=13, color=BLACK
    )

    # Card 2 — SBOM component report
    add_rect(slide, Inches(6.9), Inches(1.55), Inches(6.1), Inches(5.15), WHITE)
    add_rect(slide, Inches(6.9), Inches(1.55), Inches(6.1), Inches(0.55), TEAL)   # taller header
    add_rect(slide, Inches(6.9), Inches(2.1), Inches(0.05), Inches(4.6), TEAL)    # left accent
    add_text_box(slide, "SBOM Component Report",
                 Inches(7.04), Inches(1.62), Inches(5.82), Inches(0.42),
                 font_size=14, bold=True, color=WHITE)
    add_bullet_box(slide,
        ["Full component inventory: name, version, supplier, licence",
         "Supports single project or all active projects",
         "Collection projects produce merged reports across all children",
         "Useful for licence compliance audits",
         "Identifies components with missing or restrictive licences",
         "Landscape A4 PDF — consistent, branded format",
         "Auto-archived as a Jenkins build artefact"],
        Inches(7.08), Inches(2.18), Inches(5.76), Inches(4.4),
        font_size=13, color=BLACK
    )


def slide_example_report(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, LGREY)
    nav_bar(slide, "Example — Security Analysis Report")
    footer_bar(slide)

    # Simulated PDF report preview
    add_rect(slide, Inches(0.3), Inches(1.0), Inches(12.7), Inches(6.2),
             RGBColor(0xf8, 0xfa, 0xfc))

    # Report header
    add_rect(slide, Inches(0.3), Inches(1.0), Inches(12.7), Inches(1.25), NAVY)
    add_text_box(slide, "Security Analysis Report",
                 Inches(0.35), Inches(1.12), Inches(12.6), Inches(0.6),
                 font_size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, "customer-api  v1.8.0",
                 Inches(0.35), Inches(1.72), Inches(12.6), Inches(0.42),
                 font_size=16, color=RGBColor(0xcb, 0xd5, 0xe1), align=PP_ALIGN.CENTER)

    # Severity chips
    sev = [("CRITICAL", "1", RED), ("HIGH", "5", AMBER),
           ("MEDIUM", "9", RGBColor(0xca, 0x8a, 0x04)), ("LOW", "4", GREEN)]
    for i, (label, count, col) in enumerate(sev):
        x = Inches(1.1 + i * 2.65)
        add_rect(slide, x, Inches(2.45), Inches(2.3), Inches(0.65), col)
        add_text_box(slide, count,
                     x, Inches(2.45), Inches(2.3), Inches(0.38),
                     font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text_box(slide, label,
                     x, Inches(2.82), Inches(2.3), Inches(0.28),
                     font_size=11, color=WHITE, align=PP_ALIGN.CENTER, bold=True)

    # Risk badge
    add_rect(slide, Inches(10.55), Inches(2.45), Inches(2.2), Inches(0.65), RED)
    add_text_box(slide, "RISK: HIGH",
                 Inches(10.55), Inches(2.55), Inches(2.2), Inches(0.5),
                 font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Table header
    cols = ["CVE ID", "Component", "Version", "Severity", "CVSS", "Analysis State", "Description"]
    widths = [1.9, 2.2, 1.1, 1.1, 0.75, 1.6, 3.95]
    x_off = [0.3]
    for w in widths[:-1]:
        x_off.append(x_off[-1] + w)

    for i, (h, w) in enumerate(zip(cols, widths)):
        add_rect(slide, Inches(x_off[i]), Inches(3.32),
                 Inches(w - 0.02), Inches(0.35), NAVY)
        add_text_box(slide, h, Inches(x_off[i] + 0.04), Inches(3.35),
                     Inches(w - 0.1), Inches(0.3),
                     font_size=10, bold=True, color=WHITE)

    rows = [
        ("CVE-2021-44228", "log4j-core",          "2.14.1", "CRITICAL", "10.0", "NOT_AFFECTED",    "JNDI disabled via JVM flag — mitigated"),
        ("CVE-2022-42003", "jackson-databind",     "2.13.0", "HIGH",     "7.5",  "IN_TRIAGE",       "Resource consumption — under review"),
        ("CVE-2023-20863", "spring-expression",    "5.3.25", "HIGH",     "6.5",  "IN_TRIAGE",       "ReDoS — investigating impact"),
        ("CVE-2022-45868", "h2database",           "2.1.210","HIGH",     "7.8",  "EXPLOITABLE",     "Upgrade to 2.2.x scheduled sprint 14"),
        ("CVE-2023-34042", "spring-security-web",  "5.7.5",  "MEDIUM",   "5.3",  "FALSE_POSITIVE",  "Unused endpoint — confirmed FP"),
        ("CVE-2023-20861", "spring-expression",    "5.3.25", "MEDIUM",   "5.9",  "RESOLVED",        "Upgraded to 6.0.8 in PR #441"),
    ]
    f_colors = {"CRITICAL": RED, "HIGH": AMBER, "MEDIUM": RGBColor(0xca, 0x8a, 0x04), "LOW": GREEN}
    st_colors = {"NOT_AFFECTED": GREEN, "IN_TRIAGE": AMBER, "EXPLOITABLE": RED,
                 "RESOLVED": GREY, "FALSE_POSITIVE": GREY}
    for r, row in enumerate(rows):
        bg = WHITE if r % 2 == 0 else RGBColor(0xf1, 0xf5, 0xf9)
        for ci, (val, w) in enumerate(zip(row, widths)):
            add_rect(slide, Inches(x_off[ci]), Inches(3.72 + r * 0.4),
                     Inches(w - 0.02), Inches(0.38), bg)
            if ci == 3:
                fc = f_colors.get(val, BLACK)
            elif ci == 5:
                fc = st_colors.get(val, BLACK)
            else:
                fc = BLACK
            add_text_box(slide, val,
                         Inches(x_off[ci] + 0.04), Inches(3.75 + r * 0.4),
                         Inches(w - 0.1), Inches(0.32),
                         font_size=10, color=fc, bold=(ci in (3, 5)))


def slide_example_sbom_report(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, LGREY)
    nav_bar(slide, "Example — SBOM Component Report")
    footer_bar(slide)

    # Report header
    add_rect(slide, Inches(0.3), Inches(1.0), Inches(12.7), Inches(1.0), TEAL)
    add_text_box(slide, "SBOM Component Report",
                 Inches(0.35), Inches(1.1), Inches(12.6), Inches(0.5),
                 font_size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, "my-billing-service  v2.4.1  |  Generated 2025-04-22",
                 Inches(0.35), Inches(1.6), Inches(12.6), Inches(0.32),
                 font_size=13, color=RGBColor(0xe2, 0xe8, 0xf0), align=PP_ALIGN.CENTER)

    # Summary row
    for i, (label, val, col) in enumerate([
            ("Total Components", "187", NAVY),
            ("Direct Dependencies", "43",  TEAL),
            ("Transitive", "144", RGBColor(0x7c, 0x3a, 0xed)),
            ("Unique Licences", "12",  GREEN),
            ("Missing Licence", "3",  AMBER),
    ]):
        x = Inches(0.3 + i * 2.55)
        add_rect(slide, x, Inches(2.1), Inches(2.4), Inches(0.7), col)
        add_text_box(slide, val, x, Inches(2.1), Inches(2.4), Inches(0.38),
                     font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text_box(slide, label, x, Inches(2.47), Inches(2.4), Inches(0.28),
                     font_size=11, color=WHITE, align=PP_ALIGN.CENTER)

    # Component table
    headers = ["Component", "Version", "Supplier", "Type", "Licence", "Scope"]
    widths = [2.8, 1.4, 2.4, 1.2, 2.6, 1.6]
    x_off = [0.3]
    for w in widths[:-1]:
        x_off.append(x_off[-1] + w)

    for i, (h, w) in enumerate(zip(headers, widths)):
        add_rect(slide, Inches(x_off[i]), Inches(3.0),
                 Inches(w - 0.02), Inches(0.35), NAVY)
        add_text_box(slide, h, Inches(x_off[i] + 0.04), Inches(3.03),
                     Inches(w - 0.1), Inches(0.3),
                     font_size=11, bold=True, color=WHITE)

    rows = [
        ("spring-boot-starter-web",  "3.2.1",  "Pivotal / VMware",    "Library", "Apache-2.0",   "Required"),
        ("jackson-databind",         "2.16.1", "FasterXML",           "Library", "Apache-2.0",   "Required"),
        ("log4j-api",                "2.22.1", "Apache",              "Library", "Apache-2.0",   "Required"),
        ("hibernate-core",           "6.4.2",  "Red Hat",             "Library", "LGPL-2.1",     "Required"),
        ("commons-lang3",            "3.14.0", "Apache",              "Library", "Apache-2.0",   "Required"),
        ("bouncy-castle-provider",   "1.78",   "Legion of the BC",    "Library", "MIT",          "Optional"),
        ("netty-handler",            "4.1.107","Netty Project",       "Library", "Apache-2.0",   "Required"),
        ("unknown-lib",              "0.9.1",  "(unknown)",           "Library", "UNKNOWN",      "Required"),
    ]
    lic_colors = {"Apache-2.0": GREEN, "MIT": GREEN, "LGPL-2.1": AMBER,
                  "UNKNOWN": RED, "GPL-2.0": RED}
    for r, row in enumerate(rows):
        bg = WHITE if r % 2 == 0 else RGBColor(0xf1, 0xf5, 0xf9)
        for ci, (val, w) in enumerate(zip(row, widths)):
            add_rect(slide, Inches(x_off[ci]), Inches(3.4 + r * 0.42),
                     Inches(w - 0.02), Inches(0.4), bg)
            fc = lic_colors.get(val, BLACK) if ci == 4 else BLACK
            add_text_box(slide, val,
                         Inches(x_off[ci] + 0.04), Inches(3.43 + r * 0.42),
                         Inches(w - 0.1), Inches(0.35),
                         font_size=11, color=fc, bold=(ci == 4 and val in ("UNKNOWN", "GPL-2.0")))


def slide_pipeline(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, LGREY)
    nav_bar(slide, "End-to-End Pipeline Overview")
    footer_bar(slide)

    add_text_box(slide, "From source code commit to vulnerability insight — fully automated",
                 Inches(0.35), Inches(0.95), Inches(12.6), Inches(0.4),
                 font_size=15, italic=True, color=GREY)

    pipeline_steps = [
        (NAVY,  "Source Code",       "Developer pushes to Git\nJenkins pipeline triggered"),
        (TEAL,  "SBOM Generation",   "cdxgen or\ncyclonedx-maven-plugin\ncreates sbom.json"),
        (RGBColor(0x7c,0x3a,0xed), "Upload to DT", "uploadSBOM() posts BOM\nto Dependency-Track\nvia REST API"),
        (AMBER, "Enrichment",        "DT resolves components\nagainst NVD, OSV,\nGitHub Advisories"),
        (RED,   "Findings",          "Vulnerabilities listed\nby severity with\nCVSS scores"),
        (GREEN, "Analysis & Report", "Team triages findings\nPDF report generated\nand archived"),
    ]

    for i, (col, title, body) in enumerate(pipeline_steps):
        x = Inches(0.25 + i * 2.15)
        add_rect(slide, x, Inches(1.55), Inches(2.0), Inches(0.52), col)
        add_text_box(slide, title,
                     x + Inches(0.05), Inches(1.62), Inches(1.9), Inches(0.38),
                     font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_rect(slide, x, Inches(2.1), Inches(2.0), Inches(1.5), WHITE)
        add_text_box(slide, body,
                     x + Inches(0.07), Inches(2.18), Inches(1.87), Inches(1.35),
                     font_size=12, color=BLACK, align=PP_ALIGN.CENTER)
        # Arrow
        if i < 5:
            add_text_box(slide, "→",
                         x + Inches(2.02), Inches(2.3), Inches(0.2), Inches(0.38),
                         font_size=20, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    # Bottom detail band
    add_rect(slide, Inches(0.25), Inches(3.85), Inches(12.75), Inches(3.35), WHITE)
    add_rect(slide, Inches(0.25), Inches(3.85), Inches(12.75), Inches(0.4), NAVY)
    add_text_box(slide, "Jenkins Shared Library — Key Pipeline Steps",
                 Inches(0.4), Inches(3.9), Inches(12.5), Inches(0.32),
                 font_size=13, bold=True, color=WHITE)

    funcs = [
        ("cdxgenRepo()", "Generate SBOM for any\nlanguage repo using cdxgen"),
        ("uploadSBOM()", "Upload SBOM and create/\nupdate DT project"),
        ("createDTCollectionProject()", "Group child projects into\nan aggregate collection"),
        ("exportDTReport()", "Generate vulnerability\nPDF for one project"),
        ("exportAllDTReports()", "Generate PDF for every\nactive DT project"),
        ("exportSBOMReport()", "Generate component\nPDF for one project"),
    ]
    for i, (fn, desc) in enumerate(funcs):
        x = Inches(0.4 + i * 2.1)
        add_rect(slide, x, Inches(4.4), Inches(1.98), Inches(0.38),
                 RGBColor(0x1e, 0x29, 0x3b))
        add_text_box(slide, fn, x + Inches(0.04), Inches(4.43),
                     Inches(1.9), Inches(0.3),
                     font_size=10, bold=True,
                     color=RGBColor(0x67, 0xe8, 0xf9))
        add_text_box(slide, desc, x + Inches(0.04), Inches(4.85),
                     Inches(1.9), Inches(0.7),
                     font_size=11, color=GREY)


def slide_swagger_api(prs):
    """API Explorer & Integration Ecosystem slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, LGREY)
    nav_bar(slide, "API Explorer & Integration Ecosystem")
    footer_bar(slide)

    add_text_box(slide,
        "Everything visible in the Dependency-Track UI is automatable via its REST API — "
        "enabling deep integration with your existing toolchain.",
        Inches(0.35), Inches(0.95), Inches(12.6), Inches(0.38),
        font_size=14, italic=True, color=GREY)

    # ── Left panel: Swagger UI ──────────────────────────────────────────────
    add_rect(slide, Inches(0.3), Inches(1.48), Inches(5.9), Inches(5.75), WHITE)
    add_rect(slide, Inches(0.3), Inches(1.48), Inches(5.9), Inches(0.52), NAVY)    # taller header
    add_text_box(slide, "Swagger UI Explorer",
                 Inches(0.44), Inches(1.55), Inches(5.62), Inches(0.4),
                 font_size=15, bold=True, color=WHITE)

    add_text_box(slide,
        "Dependency-Track ships a built-in interactive OpenAPI / Swagger explorer.\n"
        "Authenticate with your API key and execute every endpoint live:",
        Inches(0.44), Inches(2.02), Inches(5.62), Inches(0.52),
        font_size=12, color=BLACK)

    # URL code-style box
    add_rect(slide, Inches(0.44), Inches(2.62), Inches(5.62), Inches(0.36),
             RGBColor(0x1e, 0x29, 0x3b))
    add_text_box(slide, "http://<dt-host>/api/swagger-ui/",
                 Inches(0.54), Inches(2.66), Inches(5.42), Inches(0.28),
                 font_size=11, bold=True, color=RGBColor(0x67, 0xe8, 0xf9))

    swagger_items = [
        "Try any endpoint interactively before writing code",
        "POST a BOM and inspect the raw API response",
        "Bulk-set analysis states via PUT /api/v1/analysis",
        "Pull live portfolio risk scores for Grafana feeds",
        "Impact analysis — find every project using a component",
        "Import dependencytrack-openapi.yaml into Postman",
    ]
    for i, item in enumerate(swagger_items):
        add_text_box(slide, f"\u2022  {item}",
                     Inches(0.54), Inches(3.1 + i * 0.52), Inches(5.52), Inches(0.45),
                     font_size=12, color=BLACK)

    # ── Right panel: Integration Ecosystem ─────────────────────────────────
    add_rect(slide, Inches(6.45), Inches(1.48), Inches(6.55), Inches(5.75), WHITE)
    add_rect(slide, Inches(6.45), Inches(1.48), Inches(6.55), Inches(0.52), TEAL)  # taller header
    add_text_box(slide, "Integration Ecosystem",
                 Inches(6.59), Inches(1.55), Inches(6.27), Inches(0.4),
                 font_size=15, bold=True, color=WHITE)

    add_text_box(slide,
        "Full visibility from developer laptop to boardroom:",
        Inches(6.59), Inches(2.02), Inches(6.27), Inches(0.32),
        font_size=12, color=BLACK)

    integrations = [
        (NAVY,  "Jira / Ticketing",
         "Auto-create issues from Critical findings;\nsync analysis state back via webhook"),
        (RED,   "Splunk / Elasticsearch",
         "Stream vulnerability events into your SIEM\nfor correlation and compliance alerting"),
        (TEAL,  "Grafana Dashboards",
         "Pull live risk scores from /metrics/portfolio\nfor real-time executive panels"),
        (AMBER, "MS Teams / Slack",
         "Webhook notifications on new Critical or\nHigh severity findings"),
        (GREEN, "PagerDuty / OpsGenie",
         "Trigger on-call alerts when newly introduced\nCritical vulnerabilities are detected"),
        (RGBColor(0x7c, 0x3a, 0xed), "Executive Reporting",
         "Scheduled metrics export — from developer\nlaptop to boardroom in one pipeline"),
    ]

    for i, (col, title, body) in enumerate(integrations):
        y = Inches(2.44 + i * 0.82)
        add_rect(slide, Inches(6.59), y, Inches(0.22), Inches(0.68), col)
        add_text_box(slide, title,
                     Inches(6.9), y + Inches(0.03), Inches(5.96), Inches(0.28),
                     font_size=12, bold=True, color=col)
        add_text_box(slide, body,
                     Inches(6.9), y + Inches(0.31), Inches(5.96), Inches(0.37),
                     font_size=11, color=GREY)


def slide_dt_ui_screenshots(prs):
    """Two-pane slide: dashboard.png top + collection-projects-details.png bottom-left
    + vulnerabilities.png bottom-right.  Skipped if no screenshots were downloaded."""
    scr_dash  = _screenshot('dashboard.png')
    scr_coll  = _screenshot('collection-projects-details.png')
    scr_vulns = _screenshot('vulnerabilities.png')

    available = [s for s in [scr_dash, scr_coll, scr_vulns] if s]
    if not available:
        return   # nothing to show — don't add an empty slide

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, LGREY)
    nav_bar(slide, "Dependency-Track — User Interface")
    footer_bar(slide)

    add_text_box(slide,
        "Real Dependency-Track UI — portfolio dashboard, project collection view, "
        "and vulnerability findings list",
        Inches(0.35), Inches(0.95), Inches(12.6), Inches(0.38),
        font_size=14, italic=True, color=GREY)

    if len(available) == 1:
        _add_screenshot_card(slide, available[0],
                             Inches(0.3), Inches(1.45), Inches(12.7), Inches(5.8),
                             "Dependency-Track Dashboard")
    elif len(available) == 2:
        titles = ["Portfolio Dashboard", "Project / Findings View"]
        for i, scr in enumerate(available):
            _add_screenshot_card(slide, scr,
                                 Inches(0.3 + i * 6.55), Inches(1.45),
                                 Inches(6.3), Inches(5.8),
                                 titles[i])
    else:
        # Top: full-width dashboard
        _add_screenshot_card(slide, available[0],
                             Inches(0.3), Inches(1.45), Inches(12.7), Inches(2.75),
                             "Portfolio Dashboard")
        # Bottom-left: collection details
        _add_screenshot_card(slide, available[1],
                             Inches(0.3), Inches(4.35), Inches(6.3), Inches(2.85),
                             "Collection Project")
        # Bottom-right: vulnerabilities
        _add_screenshot_card(slide, available[2],
                             Inches(6.73), Inches(4.35), Inches(6.27), Inches(2.85),
                             "Vulnerability Findings")


def slide_dt_audit_screenshot(prs):
    """Full-slide showcase of the real audit/analysis panel screenshot.
    Skipped if the image was not downloaded."""
    scr = _screenshot('audit-finding-project.png')
    if not scr:
        return

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, LGREY)
    nav_bar(slide, "Dependency-Track — Auditing a Finding")
    footer_bar(slide)

    add_text_box(slide,
        "The audit panel lets analysts set the analysis state, record a justification, "
        "add comments, and suppress the finding — all with a full audit trail.",
        Inches(0.35), Inches(0.95), Inches(12.6), Inches(0.38),
        font_size=14, italic=True, color=GREY)

    _add_screenshot_card(slide, scr,
                         Inches(0.3), Inches(1.45), Inches(12.7), Inches(5.8),
                         "Audit Finding — Analysis Panel")


def _add_screenshot_card(slide, img_path, left, top, width, height, caption=""):
    """Embed a screenshot image preserving aspect ratio, centred in the card area,
    with a thin navy caption bar underneath."""
    cap_h = Inches(0.33)
    img_area_w = width
    img_area_h = height - cap_h

    # Determine natural image dimensions and fit inside card area
    with _PILImage.open(img_path) as im:
        nat_w, nat_h = im.size
    scale = min(img_area_w / nat_w, img_area_h / nat_h)
    fit_w = int(nat_w * scale)
    fit_h = int(nat_h * scale)

    # Centre the fitted image within the card area
    img_left = left + (img_area_w - fit_w) // 2
    img_top  = top  + (img_area_h - fit_h) // 2

    # Fill background of card area so letterbox bars look intentional
    add_rect(slide, left, top, width, img_area_h, RGBColor(0x0a, 0x0a, 0x0a))
    slide.shapes.add_picture(img_path, img_left, img_top, fit_w, fit_h)

    add_rect(slide, left, top + img_area_h, width, cap_h, NAVY)
    if caption:
        add_text_box(slide, caption,
                     left + Inches(0.08), top + img_area_h + Inches(0.03),
                     width - Inches(0.16), cap_h - Inches(0.06),
                     font_size=11, bold=True, color=WHITE)


def slide_benefits(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, LGREY)
    nav_bar(slide, "Benefits Summary")
    footer_bar(slide)

    items = [
        (NAVY, "Complete Visibility",
         "Every dependency — direct and transitive — is known and "
         "tracked. No blind spots in your software supply chain."),
        (TEAL, "Continuous Monitoring",
         "New CVEs are automatically matched against your component "
         "inventory. You are notified without running a scan."),
        (RGBColor(0x7c, 0x3a, 0xed), "Structured Triage",
         "Analysis states, analyst comments, and suppression provide "
         "a governed workflow so decisions are documented and auditable."),
        (AMBER, "Automated Reporting",
         "Shareable PDF reports generated on demand or on a schedule — "
         "no manual work needed to produce management-ready output."),
        (RED, "Shift-Left Security",
         "In-house Maven builds upload SBOMs at build time, surfacing "
         "vulnerabilities before code reaches production."),
        (GREEN, "Policy Enforcement",
         "Define rules for severity thresholds and licence types. "
         "Violations break the build before deployment."),
    ]

    for i, (col, title, body) in enumerate(items):
        row, col_n = divmod(i, 3)
        x = Inches(0.3 + col_n * 4.35)
        y = Inches(1.3 + row * 2.85)
        add_rect(slide, x, y, Inches(4.2), Inches(2.6), WHITE)
        add_rect(slide, x, y, Inches(4.2), Inches(0.58), col)       # taller header strip
        add_rect(slide, x, y, Inches(0.05), Inches(2.6), col)        # left accent stripe
        add_text_box(slide, title,
                     x + Inches(0.18), y + Inches(0.1), Inches(3.87), Inches(0.42),
                     font_size=16, bold=True, color=WHITE)
        add_text_box(slide, body,
                     x + Inches(0.18), y + Inches(0.68), Inches(3.87), Inches(1.8),
                     font_size=13, color=BLACK, line_spacing=1.2)


# ── Assemble presentation ──────────────────────────────────────────────────

def build_pptx(out_path="SBOM_DependencyTrack_Presentation.pptx"):
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_title(prs)
    slide_what_is_sbom(prs)
    slide_how_sboms_generated(prs)
    slide_maven_plugin(prs)
    slide_dt_overview(prs)
    slide_detection_tools(prs)
    slide_dt_ui_screenshots(prs)
    slide_vuln_workflow(prs)
    slide_analysis_actions(prs)
    slide_dt_audit_screenshot(prs)
    slide_suppression(prs)
    slide_reporting(prs)
    slide_example_report(prs)
    slide_example_sbom_report(prs)
    slide_pipeline(prs)
    slide_swagger_api(prs)
    slide_benefits(prs)

    prs.save(out_path)
    print(f"Saved: {out_path}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build_pptx()
